"""Add use-case walkthrough and GRAILS safeguard slides to the federated healthcare deck.

The federated deck already explains federated learning, the round lifecycle, and the
governance argument. This script adds two things it was missing:

1. A concrete walkthrough of what happens to one patient's record, and an explicit
   list of what does and does not cross the hospital boundary.
2. A GRAILS section. GRAILS (Kulkarni & Ramanathan, AIES 2025, IIIT Bangalore) is a
   published framework for embedding ethical safeguards in software. It splits into
   Ethical-Restriction Rails (decide the protection plan) and an Ethical Guard
   (apply it), which is the same plan/enforce split Dagents already uses.

Run:
    .venv/bin/python scripts/add_grails_federated_slides.py

The script edits the deck in place and refuses to run twice. To rebuild, restore the
deck first:  git checkout -- <deck path>
"""

from __future__ import annotations

import copy
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DECK = (
    ROOT
    / "docs/presentation/healthcare-case-study/federated-healthcare-use-case"
    / "dagents-federated-healthcare-use-case.pptx"
)

# ---------------------------------------------------------------------------
# Design tokens, read off the existing slides so new ones match exactly.
# ---------------------------------------------------------------------------

BG = RGBColor(0xF7, 0xF3, 0xEA)
INK = RGBColor(0x15, 0x20, 0x1B)
MUTED = RGBColor(0x53, 0x63, 0x5B)
BORDER = RGBColor(0xB9, 0xC2, 0xBB)

GREEN = RGBColor(0xDC, 0xEB, 0xE1)
BLUE = RGBColor(0xDC, 0xE7, 0xF1)
SAND = RGBColor(0xF4, 0xE8, 0xD0)
PURPLE = RGBColor(0xE9, 0xE1, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Aptos"


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _text(slide, left, top, width, height, text, size, bold, color, anchor=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    if anchor is not None:
        tf.vertical_anchor = anchor
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, eyebrow, title, subtitle) -> None:
    _text(slide, 0.60, 0.35, 2.60, 0.32, eyebrow, 13, True, MUTED)
    _text(slide, 0.60, 0.72, 11.40, 0.65, title, 29, True, INK)
    _text(slide, 0.62, 1.34, 10.70, 0.57, subtitle, 14, False, MUTED)


def add_footer(slide, text) -> None:
    _text(slide, 0.62, 6.95, 6.20, 0.30, text, 8, False, MUTED)


def add_card(slide, left, top, width, height, title, body, fill, title_size=15, body_size=12):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    shape.text_frame.text = ""

    inner_w = width - 0.35
    title_h = 0.35 if len(title) < 26 else 0.60
    _text(slide, left + 0.18, top + 0.16, inner_w, title_h, title, title_size, True, INK)
    if body:
        body_top = top + 0.16 + title_h + 0.07
        body_h = max(0.3, height - (body_top - top) - 0.14)
        _text(slide, left + 0.18, body_top, inner_w, body_h, body, body_size, False, MUTED)
    return shape


def add_arrow(slide, left, top) -> None:
    _text(slide, left, top, 0.35, 0.57, "→", 28, True, MUTED)


def add_note(slide, text, top=5.35, size=12) -> None:
    """Single emphasised line under the main content."""
    _text(slide, 1.00, top, 10.80, 0.50, text, size, True, INK)


def add_chain(slide, steps, top=2.05, height=2.80, note=None, note_top=5.25):
    """Four (or three) cards separated by arrows, matching the round-sequence slides."""
    count = len(steps)
    if count == 4:
        left, width, gap = 0.75, 2.65, 3.05
    else:
        left, width, gap = 0.85, 3.55, 3.95
    fills = [GREEN, BLUE, SAND, PURPLE]
    for idx, (title, body) in enumerate(steps):
        x = left + idx * gap
        add_card(slide, x, top, width, height, title, body, fills[idx % 4])
        if idx < count - 1:
            add_arrow(slide, x + width + 0.05, top + height / 2 - 0.28)
    if note:
        add_note(slide, note, top=note_top)


def add_grid(slide, cards, top=1.90, height=1.25, gap_y=1.80):
    """Two-by-two card grid, matching the problem/example slides."""
    fills = [GREEN, BLUE, SAND, PURPLE]
    for idx, (title, body) in enumerate(cards):
        x = 0.85 if idx % 2 == 0 else 6.80
        y = top + (idx // 2) * gap_y
        add_card(slide, x, y, 5.35, height, title, body, fills[idx % 4])


def add_table(slide, rows, left=0.60, top=1.85, widths=(3.12, 4.44, 4.44), row_h=0.78):
    """Header row in sand, body rows in white, matching the method-boundary slides."""
    for r_idx, row in enumerate(rows):
        x = left
        for c_idx, cell in enumerate(row):
            w = widths[c_idx]
            fill = SAND if r_idx == 0 else WHITE
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(top + r_idx * row_h), Inches(w), Inches(row_h)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
            shape.line.color.rgb = BORDER
            shape.line.width = Pt(1)
            shape.shadow.inherit = False
            shape.text_frame.text = ""
            _text(
                slide,
                x + 0.08,
                top + r_idx * row_h + 0.08,
                w - 0.16,
                row_h - 0.16,
                cell,
                11 if r_idx == 0 else 10,
                r_idx == 0,
                INK if r_idx == 0 else MUTED,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            x += w


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    return slide


# ---------------------------------------------------------------------------
# Use-case depth
# ---------------------------------------------------------------------------


def follow_one_patient(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "WALKTHROUGH",
        "What happens to one patient's record.",
        "The record stays in the hospital. Only a number leaves.",
    )
    add_chain(
        slide,
        [
            (
                "ARRIVES",
                "A patient is admitted. Vital signs, laboratory results, and notes land in "
                "the hospital's own systems, exactly as they do today.",
            ),
            (
                "STAYS LOCAL",
                "The LMA builds the approved features inside the hospital. Nothing is copied "
                "out, and the hospital can refuse the job.",
            ),
            (
                "BECOMES MATH",
                "The local model trains or scores. This patient turns into part of a total, "
                "not a row anyone can read.",
            ),
            (
                "LEAVES SMALL",
                "Only the protected update and the agreed counts leave. Nobody outside can "
                "rebuild the patient from them.",
            ),
        ],
        note="The record never moves. The learning does.",
    )
    add_footer(slide, "One patient through one round / 05")
    return slide


def what_leaves(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "DATA BOUNDARY",
        "Exactly what crosses the hospital boundary.",
        "If you can name every field that leaves, you can govern it. This is the whole list.",
    )
    add_card(
        slide,
        0.85,
        1.90,
        5.35,
        3.05,
        "LEAVES THE HOSPITAL",
        "•  The protected model update\n"
        "•  Counts and scores that were agreed in advance\n"
        "•  Round number, model version, and status\n"
        "•  A pointer to evidence that stays at the hospital",
        GREEN,
    )
    add_card(
        slide,
        6.80,
        1.90,
        5.35,
        3.05,
        "NEVER LEAVES",
        "•  Patient rows\n"
        "•  Names, identifiers, and dates of birth\n"
        "•  Results for any single patient\n"
        "•  Notes and images\n"
        "•  Anything not written into the contract",
        SAND,
    )
    add_note(
        slide,
        "“We only send weights” is a claim. A named list plus a check that blocks "
        "everything else is a control.",
        top=5.30,
    )
    add_footer(slide, "Egress contract / 06")
    return slide


# ---------------------------------------------------------------------------
# GRAILS section
# ---------------------------------------------------------------------------


def grails_plain(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS",
        "A published way to decide who may see what.",
        "GRAILS keeps the ethics rules out of the normal application code, so they can be "
        "changed without rewriting the app.",
    )
    add_grid(
        slide,
        [
            (
                "Two parts",
                "The Rails work out what protection is needed. The Guard applies it. "
                "Deciding and enforcing stay separate.",
            ),
            (
                "Three questions",
                "How sensitive is the data? How far do we trust who is asking? "
                "How much are they asking for?",
            ),
            (
                "Added, not rebuilt",
                "It is designed to plug into software that already exists. "
                "You do not redesign the app to add ethics to it.",
            ),
            (
                "Wider than privacy",
                "Privacy is the worked example in the paper. The same shape is meant to "
                "extend to fairness, transparency, and security.",
            ),
        ],
    )
    add_note(
        slide,
        "Access rules based on roles alone cannot express “this person, this field, "
        "this purpose, today”. GRAILS can.",
        top=5.45,
    )
    add_footer(slide, "Kulkarni & Ramanathan, AIES 2025, IIIT Bangalore / 19")
    return slide


def grails_dials(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS",
        "Three dials decide the answer.",
        "The same request can get a different answer, depending on where the three dials sit.",
    )
    add_chain(
        slide,
        [
            (
                "HOW SENSITIVE",
                "Low, medium, or high. Set from the data itself and the rules that cover it, "
                "such as HIPAA. A ward name and a diagnosis are not the same.",
            ),
            (
                "HOW TRUSTED",
                "The paper calls this a Know Your User score, like a bank checking a customer. "
                "Built from verified identity, employer, stated purpose, and past record.",
            ),
            (
                "HOW MUCH",
                "One value, one patient, one field across everyone, or the whole table. "
                "Asking for more should not be as easy as asking for less.",
            ),
        ],
        top=2.05,
        height=2.90,
        note="High sensitivity and low trust means heavy protection. Low sensitivity and high "
        "trust means show it as it is.",
        note_top=5.35,
    )
    add_footer(slide, "Sensitivity, trust, and how much is asked for / 20")
    return slide


def grails_table(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS",
        "The dials point to a specific action.",
        "A shortened view of the paper's strategy table, written in plain terms.",
    )
    add_table(
        slide,
        [
            ["How sensitive the field is", "Asked by low trust", "Asked by medium trust", "Asked by high trust"],
            ["High", "Hide the value, or replace it with noise", "Round it off, or report it as a group", "Show it with a little noise added"],
            ["Medium", "Report it as a group only", "Show a range, hide the rest", "Show it as it is"],
            ["Low", "Add a little noise", "Show it as it is", "Show it as it is"],
        ],
        widths=(3.00, 3.20, 3.20, 3.20),
        row_h=0.80,
    )
    add_note(
        slide,
        "The point is not the exact technique. The point is that the answer is looked up from "
        "a table, not decided inside application code.",
        top=5.35,
    )
    add_footer(slide, "Filtering strategies, simplified from the paper / 21")
    return slide


def grails_fits(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS + DAGENTS",
        "GRAILS splits the same way Dagents already splits.",
        "The Rails are planning. The Guard is runtime. That line is the one Dagents is built on.",
    )
    add_card(
        slide,
        0.85,
        1.85,
        5.35,
        1.60,
        "Rails = planning",
        "Work out the protection plan. Pure, typed, repeatable, and no input or output. "
        "This belongs in the OCaml layer, beside the dataset and model planners.",
        PURPLE,
    )
    add_card(
        slide,
        6.80,
        1.85,
        5.35,
        1.60,
        "Guard = runtime",
        "Apply the plan to real data and record what was decided. This belongs in the LMA "
        "and GMA, where the side effects already are.",
        BLUE,
    )
    add_table(
        slide,
        [
            ["GRAILS part", "Where it goes in Dagents", "Why there"],
            ["Restriction Planner", "OCaml planner module", "The paper compares it to a database query planner. Dagents already has that layer."],
            ["Ethical Guard", "LMA and GMA request path", "Enforcement needs real data, real users, and an audit log."],
            ["Rules and regulations", "Configuration, not code", "Policy changes far more often than the framework does."],
        ],
        top=3.62,
        widths=(2.80, 3.40, 5.80),
        row_h=0.72,
    )
    add_footer(slide, "Plan first, then enforce / 22")
    return slide


def grails_placement(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS + DAGENTS",
        "Four places the guard has to stand.",
        "A check at the edge only is a warning label. A check at every boundary is a control.",
    )
    add_chain(
        slide,
        [
            (
                "BEFORE READING",
                "The LMA checks the job before it touches a single hospital field. "
                "An unapproved job stops here.",
            ),
            (
                "BEFORE TRAINING",
                "Check the group of patients is large enough and that every field used "
                "is on the approved list.",
            ),
            (
                "BEFORE SENDING",
                "The update and the counts are checked before anything leaves the "
                "hospital network.",
            ),
            (
                "BEFORE RELEASE",
                "The GMA checks the candidate model and its evidence before it signs "
                "anything as released.",
            ),
        ],
        note="Every check writes down what it decided, why, and which round it belonged to.",
    )
    add_footer(slide, "Where the guard sits / 23")
    return slide


def grails_gap(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS + FEDERATION",
        "GRAILS protects rows. Federation sends updates.",
        "This is the one place the framework has to be extended for this use case, and it is "
        "the contribution this project can make.",
    )
    add_card(
        slide,
        0.85,
        1.85,
        5.35,
        1.45,
        "What it covers today",
        "One value, one row, one column, or a whole table. All of them are data you can "
        "point at and read.",
        GREEN,
    )
    add_card(
        slide,
        6.80,
        1.85,
        5.35,
        1.45,
        "What federation adds",
        "A model update is not a row, but it still carries patient signal. It needs to become "
        "a fifth thing the rules cover.",
        SAND,
    )
    add_card(
        slide,
        0.85,
        3.55,
        3.75,
        1.55,
        "Sensitivity becomes",
        "How much patient signal an update can carry out of the hospital.",
        PURPLE,
    )
    add_card(
        slide,
        4.85,
        3.55,
        3.75,
        1.55,
        "Trust becomes",
        "The receiving site and the coordinator, not only a named person.",
        BLUE,
    )
    add_card(
        slide,
        8.85,
        3.55,
        3.30,
        1.55,
        "Action becomes",
        "Limit the contribution, require enough sites, add noise, or refuse to send.",
        GREEN,
    )
    add_note(
        slide,
        "This is a real extension, not a restatement of the paper.",
        top=5.35,
    )
    add_footer(slide, "Extending the rules to model updates / 24")
    return slide


def grails_example(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS",
        "One request, from start to finish.",
        "A researcher at one hospital asks for stroke cases so they can check a model.",
    )
    add_chain(
        slide,
        [
            (
                "WHO IS ASKING",
                "A verified staff account, a named study, an approved purpose, and a clean "
                "record. Trust comes out as medium.",
            ),
            (
                "WHAT THEY WANT",
                "Arrival time, age, and outcome for stroke patients. That is one field across "
                "many patients, not one patient.",
            ),
            (
                "WHAT THE RULES SAY",
                "Outcome is a high sensitivity field. Medium trust, asked for across a whole "
                "column, so group it and hide part of it.",
            ),
            (
                "WHAT THEY GET",
                "Age bands instead of ages, grouped outcomes, and a written reason. Enough to "
                "check the model, not enough to identify anyone.",
            ),
        ],
        note="Nobody wrote code for this specific request. The plan came from the three dials.",
    )
    add_footer(slide, "A worked request / 25")
    return slide


# ---------------------------------------------------------------------------
# Ordering and renumbering
# ---------------------------------------------------------------------------


def move_slide(prs, slide, new_index) -> None:
    sld_id_lst = prs.slides._sldIdLst
    target = None
    for element in sld_id_lst:
        if element.get("id") == str(slide.slide_id):
            target = element
            break
    if target is None:
        raise RuntimeError(f"slide id {slide.slide_id} not found")
    sld_id_lst.remove(target)
    sld_id_lst.insert(new_index, target)


FOOTER_NUM = re.compile(r"(/\s*)(\d+)\s*$")


def renumber_footers(prs) -> int:
    """Rewrite the trailing '/ NN' in each footer so inserted slides stay correct."""
    fixed = 0
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if abs(shape.top / 914400 - 6.95) > 0.12:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if FOOTER_NUM.search(run.text):
                        new = FOOTER_NUM.sub(lambda m: f"{m.group(1)}{idx:02d}", run.text)
                        if new != run.text:
                            run.text = new
                            fixed += 1
    return fixed


def main() -> None:
    prs = Presentation(str(DECK))
    before = len(prs.slides._sldIdLst)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == "GRAILS":
                raise SystemExit(
                    "Deck already contains GRAILS slides. Restore it first:\n"
                    "  git checkout -- " + str(DECK.relative_to(ROOT))
                )

    # Deepen the use case: goes right after the use-case boundary slide (4).
    use_case = [follow_one_patient(prs), what_leaves(prs)]
    # GRAILS section: goes after the round safety gates slide, before the options slide.
    grails = [
        grails_plain(prs),
        grails_dials(prs),
        grails_table(prs),
        grails_fits(prs),
        grails_placement(prs),
        grails_gap(prs),
        grails_example(prs),
    ]

    for offset, slide in enumerate(use_case):
        move_slide(prs, slide, 4 + offset)          # after original slide 4
    for offset, slide in enumerate(grails):
        move_slide(prs, slide, 18 + offset)         # after gates slide, now index 18

    fixed = renumber_footers(prs)
    prs.save(str(DECK))
    after = len(Presentation(str(DECK)).slides._sldIdLst)
    print(f"slides {before} -> {after}; footers renumbered: {fixed}")


if __name__ == "__main__":
    main()
