"""Add a stroke-specific GRAILS section to the federated stroke-triage deck.

The companion script `add_grails_federated_slides.py` adds a general GRAILS section to
the sepsis-framed federated deck. This one is written for the stroke deck, so the
examples are stroke fields and stroke requesters rather than generic ones, and the
guard is mapped onto the four round safety gates that deck already presents.

GRAILS: Kulkarni & Ramanathan, "GRAILS - A Framework for Embedding Ethical Safeguards
in Software Applications for Responsible AI", AIES 2025, IIIT Bangalore.

Run:
    .venv/bin/python scripts/add_grails_stroke_slides.py

The script edits the deck in place and refuses to run twice. To rebuild, restore the
deck first:  git checkout -- <deck path>
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation

from add_grails_federated_slides import (  # shared design tokens and layout helpers
    BLUE,
    GREEN,
    PURPLE,
    SAND,
    add_card,
    add_chain,
    add_footer,
    add_grid,
    add_note,
    add_table,
    add_title,
    move_slide,
    new_slide,
    renumber_footers,
)

ROOT = Path(__file__).resolve().parents[1]
DECK = (
    ROOT
    / "docs/presentation/healthcare-case-study/federated-healthcare-use-case"
    / "dagents-federated-stroke-use-case.pptx"
)

# The GRAILS section is inserted directly after the round safety gates slide, so the
# argument runs: here are the gates -> here is what makes gate three concrete -> here
# are the options and the roadmap.
INSERT_AFTER_TITLE = "Four gates must pass before a round changes the model."


def grails_intro(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS",
        "A published way to decide who may see what.",
        "GRAILS keeps the ethics rules out of the normal application code, so a policy "
        "change does not become a code change.",
    )
    add_grid(
        slide,
        [
            (
                "Two parts",
                "The Rails work out what protection a request needs. The Guard applies it. "
                "Deciding and enforcing stay separate.",
            ),
            (
                "Three questions",
                "How sensitive is the data? How far do we trust who is asking? "
                "How much are they asking for?",
            ),
            (
                "Added, not rebuilt",
                "It plugs into software that already exists. The stroke pathway is not "
                "redesigned to add ethics to it.",
            ),
            (
                "Wider than privacy",
                "Privacy is the paper's worked example. The same shape extends to "
                "fairness, transparency, and security.",
            ),
        ],
    )
    add_note(
        slide,
        "Access rules based on job titles alone cannot express “this person, this field, "
        "this purpose, today”. GRAILS can.",
        top=5.45,
    )
    add_footer(slide, "Kulkarni & Ramanathan, AIES 2025, IIIT Bangalore / 21")
    return slide


def grails_dials_stroke(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS + STROKE",
        "The three dials, using stroke data.",
        "The same request can get a different answer, depending on where the three dials sit.",
    )
    add_chain(
        slide,
        [
            (
                "HOW SENSITIVE",
                "A ward name is low. The time symptoms started is medium. The brain scan and "
                "the final diagnosis are high.",
            ),
            (
                "HOW TRUSTED",
                "A treating clinician on the stroke pathway, an internal audit reviewer, and "
                "an outside researcher are not the same requester.",
            ),
            (
                "HOW MUCH",
                "One value for one patient, one patient's whole record, one field across "
                "every stroke case, or the entire table.",
            ),
        ],
        top=2.05,
        height=2.90,
        note="High sensitivity, asked for by low trust, across a whole column is the case that "
        "must never quietly default to allow.",
        note_top=5.35,
    )
    add_footer(slide, "Sensitivity, trust, and how much is asked for / 22")
    return slide


def grails_stroke_fields(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS + STROKE",
        "The same field, three different answers.",
        "What a requester gets back depends on who they are, not only on what they asked for.",
    )
    add_table(
        slide,
        [
            ["Stroke field", "Treating clinician", "Audit reviewer", "Outside researcher"],
            ["When symptoms started", "Exact time", "Rounded to the hour", "Time band only"],
            ["Brain scan image", "Full image", "Report text, not the image", "Not released"],
            ["Stroke severity score", "Exact score", "Grouped into bands", "Bands, small groups hidden"],
            ["Final diagnosis", "Full detail", "Grouped category", "Counts only, minimum group size"],
        ],
        top=1.82,
        widths=(3.10, 3.00, 3.20, 3.30),
        row_h=0.74,
    )
    add_note(
        slide,
        "None of this is an if-statement inside the application. The row is looked up from "
        "the three dials.",
        top=5.62,
    )
    add_footer(slide, "Stroke fields under the three dials / 23")
    return slide


def grails_fits(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS + DAGENTS",
        "GRAILS splits the same way Dagents already splits.",
        "The Rails are planning. The Guard is runtime. That line is the one Dagents is built on.",
    )
    add_card(
        slide,
        0.85,
        1.85,
        5.35,
        1.60,
        "Rails = planning",
        "Work out the protection plan. Pure, typed, repeatable, and no input or output. "
        "This belongs in the OCaml layer, beside the dataset and model planners.",
        PURPLE,
    )
    add_card(
        slide,
        6.80,
        1.85,
        5.35,
        1.60,
        "Guard = runtime",
        "Apply the plan to real data and record what was decided. This belongs in the LMA "
        "and GMA, where the side effects already are.",
        BLUE,
    )
    add_table(
        slide,
        [
            ["GRAILS part", "Where it goes in Dagents", "Why there"],
            [
                "Restriction Planner",
                "OCaml planner module",
                "The paper compares it to a database query planner. Dagents already has that layer.",
            ],
            [
                "Ethical Guard",
                "LMA and GMA request path",
                "Enforcement needs real data, a real requester, and an audit log.",
            ],
            [
                "Rules and regulations",
                "Configuration, not code",
                "Hospital policy changes far more often than the framework does.",
            ],
        ],
        top=3.62,
        widths=(2.80, 3.40, 5.80),
        row_h=0.72,
    )
    add_footer(slide, "Plan first, then enforce / 24")
    return slide


def grails_on_gates(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS + STROKE",
        "The guard stands at the four gates you already have.",
        "GRAILS does not add a fifth gate. It makes the existing ones concrete and reviewable.",
    )
    add_chain(
        slide,
        [
            (
                "1. SCOPE + SITE",
                "The trust score decides whether this hospital, this study, and this "
                "requester may run at all.",
            ),
            (
                "2. CODE + DATA",
                "Only fields on the stroke feature contract can be read. Anything else is "
                "refused before a scan is ever opened.",
            ),
            (
                "3. PRIVACY",
                "The restriction plan sets the limit on the update and the counts before "
                "they leave the hospital network.",
            ),
            (
                "4. RELEASE",
                "The GMA checks the plan was followed and the evidence exists before it "
                "signs anything as released.",
            ),
        ],
        note="Every decision writes down what it decided, why, and which round it belonged to.",
    )
    add_footer(slide, "Guard placement on the existing gates / 25")
    return slide


def grails_gap(prs):
    slide = new_slide(prs)
    add_title(
        slide,
        "GRAILS + FEDERATION",
        "GRAILS protects rows. A stroke round sends updates.",
        "This is the one place the framework has to be extended, and it is the contribution "
        "this project can make.",
    )
    add_card(
        slide,
        0.85,
        1.85,
        5.35,
        1.45,
        "What it covers today",
        "One value, one row, one column, or a whole table. All of them are data you can "
        "point at and read.",
        GREEN,
    )
    add_card(
        slide,
        6.80,
        1.85,
        5.35,
        1.45,
        "What a stroke round adds",
        "A model update is not a row, but it still carries patient signal. It has to become "
        "a fifth thing the rules cover.",
        SAND,
    )
    add_card(
        slide,
        0.85,
        3.55,
        3.75,
        1.55,
        "Sensitivity becomes",
        "How much stroke-patient signal an update can carry out of the hospital.",
        PURPLE,
    )
    add_card(
        slide,
        4.85,
        3.55,
        3.75,
        1.55,
        "Trust becomes",
        "The receiving site and the coordinator, not only a named person.",
        BLUE,
    )
    add_card(
        slide,
        8.85,
        3.55,
        3.30,
        1.55,
        "Action becomes",
        "Limit the contribution, require enough sites, add noise, or refuse to send.",
        GREEN,
    )
    add_note(
        slide,
        "This is a real extension of the published framework, not a restatement of it.",
        top=5.35,
    )
    add_footer(slide, "Extending the rules to model updates / 26")
    return slide


def main() -> None:
    prs = Presentation(str(DECK))
    before = len(prs.slides._sldIdLst)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == "GRAILS":
                raise SystemExit(
                    "Deck already contains GRAILS slides. Restore it first:\n"
                    "  git checkout -- " + str(DECK.relative_to(ROOT))
                )

    anchor = None
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == INSERT_AFTER_TITLE:
                anchor = idx
                break
        if anchor is not None:
            break
    if anchor is None:
        raise SystemExit(f"could not find the anchor slide: {INSERT_AFTER_TITLE!r}")

    section = [
        grails_intro(prs),
        grails_dials_stroke(prs),
        grails_stroke_fields(prs),
        grails_fits(prs),
        grails_on_gates(prs),
        grails_gap(prs),
    ]
    for offset, slide in enumerate(section):
        move_slide(prs, slide, anchor + 1 + offset)

    fixed = renumber_footers(prs)
    prs.save(str(DECK))
    after = len(Presentation(str(DECK)).slides._sldIdLst)
    print(f"anchor slide index {anchor + 1}; slides {before} -> {after}; footers renumbered: {fixed}")


if __name__ == "__main__":
    main()
