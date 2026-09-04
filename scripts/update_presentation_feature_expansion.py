from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "docs/presentation/dagents-project-presentation.pptx"
PUML = ROOT / "docs/presentation/puml"


BG = RGBColor(247, 243, 234)
INK = RGBColor(21, 32, 27)
MUTED = RGBColor(83, 99, 91)
GREEN = RGBColor(220, 235, 225)
BLUE = RGBColor(220, 231, 241)
TAN = RGBColor(244, 232, 208)
PURPLE = RGBColor(233, 225, 250)
RED = RGBColor(240, 209, 200)
WHITE = RGBColor(255, 255, 255)


def set_bg(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 18,
             bold: bool = False, color: RGBColor = INK, align=PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"


def add_title(slide, section: str, headline: str, subhead: str | None = None) -> None:
    add_text(slide, section.upper(), 0.6, 0.35, 2.6, 0.3, size=13, bold=True, color=MUTED)
    add_text(slide, headline, 0.6, 0.72, 11.4, 0.65, size=29, bold=True)
    if subhead:
        add_text(slide, subhead, 0.62, 1.34, 10.7, 0.45, size=14, color=MUTED)


def add_footer(slide, label: str) -> None:
    add_text(slide, label, 0.62, 6.95, 4.0, 0.25, size=8, color=MUTED)


def add_card(slide, x: float, y: float, w: float, h: float, title: str, body: str,
             fill: RGBColor = WHITE) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(185, 194, 187)
    shape.line.width = Pt(1)
    add_text(slide, title, x + 0.18, y + 0.16, w - 0.35, 0.3, size=15, bold=True)
    add_text(slide, body, x + 0.18, y + 0.58, w - 0.35, h - 0.72, size=12, color=MUTED)


def add_bullet_block(slide, items: list[str], x: float, y: float, w: float, h: float,
                     size: int = 16) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Aptos"
        p.font.color.rgb = INK


def add_image(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_table_like(slide, rows: list[tuple[str, str, str]], x: float, y: float, w: float, h: float) -> None:
    row_h = h / len(rows)
    for i, (a, b, c) in enumerate(rows):
        fy = y + i * row_h
        fill = TAN if i == 0 else WHITE
        for j, (text, cw) in enumerate([(a, 0.26), (b, 0.37), (c, 0.37)]):
            fx = x + w * sum([0.26, 0.37, 0.37][:j])
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(fx), Inches(fy), Inches(w * cw), Inches(row_h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
            shape.line.color.rgb = RGBColor(190, 190, 185)
            add_text(slide, text, fx + 0.08, fy + 0.08, w * cw - 0.16, row_h - 0.12, size=10 if i else 11, bold=(i == 0))


def add_section_divider(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_text(slide, title, 0.75, 2.25, 11.0, 0.75, size=36, bold=True)
    add_text(slide, subtitle, 0.78, 3.05, 9.8, 0.55, size=18, color=MUTED)
    add_footer(slide, "Dagents / expanded progress")
    return slide


def progress_snapshot(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "Progress", "What is built now", "Dagents is now more than an idea: the repo has running services, agents, planners, a demo app, and scripts.")
    cards = [
        ("Agents", "LMA handles source-level profiling and model runs. GMA handles registration, aggregate profiling, and dispatch.", GREEN),
        ("Services", "Core, pipeline, and model services expose real APIs for catalog, topology, workflows, and ML jobs.", TAN),
        ("Functional planners", "OCaml planners validate source specs, schema contracts, quality rules, pipeline DAGs, model routes, and manifests.", PURPLE),
        ("Demo app", "NL2SQL calls Dagents during a normal /generate request and returns a trace showing what happened.", BLUE),
        ("Deployment", "Docker Compose starts the full local stack. Scripts build, probe, and stop the demo.", WHITE),
        ("Docs", "Slides, diagrams, final report, debug notes, and expected demo output are now in the repo.", WHITE),
    ]
    for idx, (title, body, fill) in enumerate(cards):
        add_card(slide, 0.65 + (idx % 3) * 4.05, 2.0 + (idx // 3) * 1.85, 3.65, 1.35, title, body, fill)
    add_footer(slide, "New section: implementation progress")
    return slide


def capability_map(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "Capabilities", "Current capability map", "The project now has reusable pieces that consumer backends can call instead of rebuilding agent infrastructure.")
    add_image(slide, PUML / "dagents_capability_map.png", 0.55, 1.65, 12.2, 5.05)
    add_footer(slide, "PlantUML: docs/presentation/puml/dagents_capability_map.puml")
    return slide


def agent_control_plane(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "Agents", "Agent control plane is working", "The LMA/GMA split gives the framework a clear local-vs-global boundary.")
    steps = [
        ("1", "Register", "GMA records LMA identity, scope, version, and capabilities."),
        ("2", "Profile", "LMA profiles source records; GMA profiles assimilated records."),
        ("3", "Run models", "Both agents can accept model-job requests and produce run metadata."),
        ("4", "Sync", "GMA can plan desired deployments and receive LMA heartbeat/sync state."),
    ]
    for i, (_, title, body) in enumerate(steps):
        add_card(slide, 0.75 + i * 3.05, 2.05, 2.65, 2.8, title, body, [GREEN, BLUE, TAN, PURPLE][i])
        if i < 3:
            add_text(slide, "→", 3.45 + i * 3.05, 3.2, 0.35, 0.35, size=28, bold=True, color=MUTED)
    add_bullet_block(
        slide,
        [
            "This is useful for any app that wants local computation per source and a combined multi-source view.",
            "The same pattern can later support Watchdog, Datalytics, and other backend systems.",
        ],
        1.0,
        5.25,
        10.8,
        0.9,
        size=15,
    )
    add_footer(slide, "Implemented agent APIs: LMA + GMA")
    return slide


def service_suite(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "Services", "The service suite is ready for integration", "Each service owns a narrow job, so apps can adopt Dagents gradually.")
    rows = [
        ("Service", "What it does", "Why it matters"),
        ("core-service", "Catalog, topology, workload compile, Kubernetes YAML", "Backend teams get deployment plans without hand-writing manifests."),
        ("pipeline-service", "Registers JSON workflows and runs ordered steps", "Workflows become inspectable and repeatable."),
        ("model-service", "Dataset catalog, training jobs, model-job visibility", "Common ML tasks move into a shared service."),
        ("LMA / GMA", "Source-level and aggregate agent computation", "One-source and multi-source jobs have separate control boundaries."),
        ("dagentsc", "OCaml process boundary for deterministic planning", "Services get typed checks without embedding OCaml directly."),
    ]
    add_table_like(slide, rows, 0.6, 1.85, 12.0, 4.65)
    add_footer(slide, "Integration is API-based, not a codebase fork")
    return slide


def nl2sql_status(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "NL2SQL Demo", "Demo app progress", "The NL2SQL app now demonstrates the full framework path, not only SQL generation.")
    add_card(slide, 0.7, 1.9, 3.7, 1.55, "Dagents trace", "The backend returns planner checks, service calls, LMA/GMA actions, and workload compilation evidence.", BLUE)
    add_card(slide, 4.8, 1.9, 3.7, 1.55, "Model adapter", "CodeT5 model artifacts are discovered from models/*.zip and loaded by the adapter when dependencies are present.", GREEN)
    add_card(slide, 8.9, 1.9, 3.7, 1.55, "Demo scripts", "run_compose_demo.sh, probe_api.sh, and stop_compose_demo.sh make the presentation path repeatable.", TAN)
    add_bullet_block(
        slide,
        [
            "The latest probe validates 22 Dagents trace steps.",
            "The model path now runs with used_fallback = false when optional model dependencies are installed.",
            "The deterministic fallback still exists for low-resource machines, but the probe can fail if fallback is not expected.",
        ],
        0.95,
        4.05,
        11.5,
        1.25,
        size=16,
    )
    add_footer(slide, "Demo output: docs/demo/expected/nl2sql-generate.json")
    return slide


def request_flow(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "Demo Flow", "What happens during /generate", "The app request becomes a traceable sequence of planner calls, service calls, and model generation.")
    add_image(slide, PUML / "dagents_request_flow.png", 0.55, 1.75, 12.1, 4.35)
    add_footer(slide, "PlantUML: docs/presentation/puml/dagents_request_flow.puml")
    return slide


def development_challenges(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "Development", "Challenges and steering decisions", "We used the AI assistant as an implementation partner, but kept decisions grounded in test results and logs.")
    rows = [
        ("Challenge", "How we steered it", "State change"),
        ("Docker model deps", "Use logs, identify CUDA wheel blow-up, pin CPU-friendly Torch line.", "Backend image now builds with model deps."),
        ("Fallback SQL", "Trace the exception instead of hiding it.", "Model adapter now loads CodeT5 artifact."),
        ("JSON key casing", "Compare trace payloads and find false schema failures.", "Record field names are preserved."),
        ("OCaml readability", "Ask for modular common_ir without changing behavior.", "Types are grouped while old API remains."),
        ("Demo reliability", "Make scripts executable and probe real service behavior.", "Probe validates trace and model path."),
    ]
    add_table_like(slide, rows, 0.55, 1.78, 12.15, 4.95)
    add_footer(slide, "This slide also answers the development-challenges report request")
    return slide


def project_potential(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "Potential", "Where this project can grow", "Dagents can become the shared agent layer for data-heavy applications.")
    items = [
        ("Reusable agent layer", "Apps can reuse LMA/GMA behavior instead of writing one-off orchestration code."),
        ("Traceable decisions", "Every request can show checks, plans, routes, and service calls."),
        ("Deployment bridge", "Typed workload specs can become Docker, Kubernetes, or cloud deployment artifacts."),
        ("Safer AI workflows", "Policy checks can run before data access, model calls, and generated actions."),
    ]
    for i, (title, body) in enumerate(items):
        add_card(slide, 0.85 + (i % 2) * 5.95, 1.9 + (i // 2) * 1.8, 5.35, 1.25, title, body, [GREEN, BLUE, TAN, PURPLE][i])
    add_footer(slide, "Potential: shared infra for future consumers")
    return slide


def future_integrations(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "Future", "Cloud and backend integrations", "The framework is service-based, so it can connect to cloud services without changing the core agent model.")
    add_image(slide, PUML / "dagents_cloud_integrations.png", 0.85, 1.75, 11.45, 4.75)
    add_footer(slide, "Examples: GCP, AWS, Supabase, Azure")
    return slide


def future_backlog(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "Future", "Concrete integration backlog", "These are practical next steps after the current local/demo stack.")
    cards = [
        ("GCP", "GKE for workloads, Cloud SQL for sources, GCS for artifacts, Pub/Sub for events.", BLUE),
        ("AWS", "EKS for workloads, S3 for artifacts, RDS for sources, SQS/EventBridge for messaging.", TAN),
        ("Supabase", "Postgres and Auth for quick product integrations; Edge Functions for lightweight hooks.", GREEN),
        ("Observability", "OpenTelemetry traces, run metrics, model-job metrics, and audit logs.", PURPLE),
        ("Secrets", "Vault/KMS/IAM integration so source credentials are never hard-coded.", WHITE),
        ("Persistence", "Move in-memory registries to Postgres or a managed store.", WHITE),
    ]
    for idx, (title, body, fill) in enumerate(cards):
        add_card(slide, 0.65 + (idx % 3) * 4.05, 1.82 + (idx // 3) * 1.75, 3.65, 1.25, title, body, fill)
    add_footer(slide, "Future integrations should keep the same LMA/GMA contracts")
    return slide


def grails_intro(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "Future", "Ethical firewalls for AI-era systems", "A GRAILS-style layer can be treated as a policy firewall around agent actions.")
    add_card(slide, 0.8, 1.85, 3.6, 2.25, "What it checks", "Should this request access this data? Should this model be called? Should this output be released?", PURPLE)
    add_card(slide, 4.85, 1.85, 3.6, 2.25, "Where it runs", "Before source extraction, before model routing, before workload deployment, and before returning generated output.", BLUE)
    add_card(slide, 8.9, 1.85, 3.6, 2.25, "What it records", "The policy decision, reason, request id, agent id, data scope, and final action.", GREEN)
    add_bullet_block(
        slide,
        [
            "This would make Dagents useful not only for computation, but also for responsible control of computation.",
            "The idea fits the LMA/GMA model because local and global decisions can have different policy rules.",
        ],
        0.95,
        4.65,
        11.2,
        1.0,
        size=16,
    )
    add_footer(slide, "Future direction informed by the GRAILS responsible-AI safeguards paper and the provided talk.")
    return slide


def grails_placement(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "Future", "Where a GRAILS-style guardrail fits", "The guardrail should be a first-class service boundary, not a UI-only warning.")
    add_image(slide, PUML / "dagents_grails_firewall.png", 0.55, 2.0, 12.2, 2.45)
    add_bullet_block(
        slide,
        [
            "Policy can stop a run, narrow the data scope, or require human review.",
            "Audit logs should connect policy decisions with Dagents run ids and deployment plans.",
            "This can support college research around ethical firewalls while staying practical for backend teams.",
        ],
        0.85,
        4.95,
        11.6,
        1.1,
        size=16,
    )
    add_footer(slide, "GRAILS integration is proposed as future work; policy mapping needs project-specific design.")
    return slide


def roadmap(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide)
    add_title(slide, "Roadmap", "Next build steps", "The path forward is incremental: harden the current stack, then add cloud and policy boundaries.")
    rows = [
        ("Phase", "Goal", "Output"),
        ("1. Stabilize", "Persist registrations and runs; improve error reporting.", "Reliable local demo and repeatable tests."),
        ("2. Deploy", "Move workloads from Compose to Kubernetes.", "Core-service manifests applied to a real cluster."),
        ("3. Integrate", "Connect GCP/AWS/Supabase data and auth services.", "Reusable adapters and secret handling."),
        ("4. Govern", "Add GRAILS-style policy checks and audit logs.", "Ethical firewall for data, models, and outputs."),
        ("5. Scale", "Broker-backed messaging and distributed agent execution.", "Production-ready multi-source computation."),
    ]
    add_table_like(slide, rows, 0.65, 1.85, 12.0, 4.45)
    add_footer(slide, "Future work: move from demo stack to production platform")
    return slide


def slide_id_element(prs: Presentation, slide) -> object:
    target = str(slide.slide_id)
    for element in prs.slides._sldIdLst:
        if element.get("id") == target:
            return element
    raise ValueError("slide id not found")


def move_slides_to_index(prs: Presentation, slides: list, index: int) -> None:
    sld_ids = prs.slides._sldIdLst
    elements = [slide_id_element(prs, slide) for slide in slides]
    for element in elements:
        sld_ids.remove(element)
    for offset, element in enumerate(elements):
        sld_ids.insert(index + offset, element)


def move_slides_before(prs: Presentation, slides: list, before_slide) -> None:
    sld_ids = prs.slides._sldIdLst
    before = slide_id_element(prs, before_slide)
    elements = [slide_id_element(prs, slide) for slide in slides]
    for element in elements:
        sld_ids.remove(element)
    index = list(sld_ids).index(before)
    for offset, element in enumerate(elements):
        sld_ids.insert(index + offset, element)


def main() -> None:
    prs = Presentation(PPTX)
    close_slide = prs.slides[-1]

    progress_slides = [
        add_section_divider(prs, "More Than OCaml", "What else is already built in Dagents"),
        progress_snapshot(prs),
        capability_map(prs),
        agent_control_plane(prs),
        service_suite(prs),
        nl2sql_status(prs),
        request_flow(prs),
        development_challenges(prs),
        project_potential(prs),
    ]

    future_slides = [
        future_integrations(prs),
        future_backlog(prs),
        grails_intro(prs),
        grails_placement(prs),
        roadmap(prs),
    ]

    move_slides_to_index(prs, progress_slides, 15)
    move_slides_before(prs, future_slides, close_slide)
    prs.save(PPTX)


if __name__ == "__main__":
    main()
